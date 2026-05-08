"""PPTX extractor: one chunk per slide.

Body = title + all text-shape contents + speaker notes. Slide number is 1-based.
Title field is populated from the slide's title placeholder when present.
"""

from __future__ import annotations

import hashlib
from collections.abc import Iterator
from pathlib import Path

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from acorn.extract.base import Block, Chunk


def _parent_id(path: Path) -> str:
    return hashlib.sha1(str(path.resolve()).encode("utf-8")).hexdigest()


def _shape_text(shape: BaseShape) -> str:
    if not shape.has_text_frame:
        return ""
    # has_text_frame == True implies the shape has a .text_frame attribute, but
    # this isn't expressed in python-pptx's static types. Pull via getattr.
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return ""
    return getattr(tf, "text", "") or ""


def _slide_title(slide: Slide) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is None:
        return ""
    return _shape_text(title_shape).strip()


def _slide_notes(slide: Slide) -> str:
    if not slide.has_notes_slide:
        return ""
    notes = slide.notes_slide
    tf = getattr(notes, "notes_text_frame", None)
    if tf is None:
        return ""
    return getattr(tf, "text", "") or ""


def extract(path: Path) -> Iterator[Chunk]:
    parent_id = _parent_id(path)
    mtime = int(path.stat().st_mtime)
    prs = Presentation(str(path))

    deck_title = ""
    for slide_index, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        if slide_index == 1 and title:
            deck_title = title

        body_parts: list[str] = []
        blocks: list[Block] = []
        if title:
            blocks.append(Block(kind="h1", text=title))
            body_parts.append(title)

        for shape in slide.shapes:
            if shape == getattr(slide.shapes, "title", None):
                continue
            text = _shape_text(shape).strip()
            if not text:
                continue
            blocks.append(Block(kind="p", text=text))
            body_parts.append(text)

        notes = _slide_notes(slide).strip()
        if notes:
            blocks.append(Block(kind="quote", text=notes))
            body_parts.append(notes)

        body = "\n".join(body_parts).strip()
        if not body:
            continue

        yield Chunk(
            parent_id=parent_id,
            path=str(path),
            mtime=mtime,
            kind="pptx",
            body=body,
            body_struct=blocks,
            slide=slide_index,
            heading_path=title,
            title=deck_title,
            chunk_seq=slide_index - 1,
        )
