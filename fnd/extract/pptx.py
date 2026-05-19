"""PPTX extractor: one chunk per slide.

Body = title + all text-shape contents + tables + speaker notes.
``body_md`` carries the slide as markdown so the structural preview
renderer (Textual Markdown widget) can render headings, bulleted
lists (paragraphs at ``level > 0``), pipe tables, and a trailing
blockquote for speaker notes. ``body_struct`` keeps the legacy
plain-text Block list for the snippet pipeline.
"""

from __future__ import annotations

import hashlib
from collections.abc import Iterator
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.exc import PackageNotFoundError
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from fnd.extract._ooxml import reject_if_zip_bomb
from fnd.extract.base import Block, Chunk, ExtractError


def _parent_id(path: Path) -> str:
    return hashlib.sha1(str(path.resolve()).encode("utf-8"), usedforsecurity=False).hexdigest()


def _shape_text(shape: BaseShape) -> str:
    if not shape.has_text_frame:
        return ""
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return ""
    return getattr(tf, "text", "") or ""


def _slide_title(slide: Slide) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is None:
        return ""
    return _shape_text(title_shape).strip()


def _slide_notes(slide: Slide) -> str:
    if not slide.has_notes_slide:
        return ""
    notes = slide.notes_slide
    tf = getattr(notes, "notes_text_frame", None)
    if tf is None:
        return ""
    return getattr(tf, "text", "") or ""


def _paragraph_md(paragraph: Any) -> str:
    """Render a python-pptx Paragraph's runs with bold/italic markers.

    Whitespace-only runs pass through bare to avoid stray ``** **``
    artefacts; bold+italic combine as ``***…***``.
    """
    parts: list[str] = []
    for run in paragraph.runs:
        text = run.text or ""
        if not text:
            continue
        bold = bool(getattr(run.font, "bold", False))
        italic = bool(getattr(run.font, "italic", False))
        if not text.strip():
            parts.append(text)
            continue
        if bold and italic:
            parts.append(f"***{text}***")
        elif bold:
            parts.append(f"**{text}**")
        elif italic:
            parts.append(f"*{text}*")
        else:
            parts.append(text)
    return "".join(parts)


def _shape_md_lines(shape: BaseShape) -> list[str]:
    """Convert a text shape to markdown lines.

    A shape with any paragraph at ``level > 0`` is treated as a
    bulleted list (every paragraph in the shape becomes a list item,
    indented by its level). Otherwise paragraphs render as separate
    paragraphs. This matches typical PowerPoint authorship: body
    placeholders set non-zero levels for sub-bullets, plain text
    boxes leave everything at level 0.
    """
    if not shape.has_text_frame:
        return []
    tf = shape.text_frame  # type: ignore[attr-defined]
    paragraphs = list(tf.paragraphs)
    if not paragraphs:
        return []
    is_bullet_list = any(p.level > 0 for p in paragraphs)
    lines: list[str] = []
    for p in paragraphs:
        rendered = _paragraph_md(p)
        if not rendered.strip():
            continue
        if is_bullet_list:
            indent = "  " * max(p.level, 0)
            lines.append(f"{indent}- {rendered}")
        else:
            lines.append(rendered)
    return lines


def _table_md(shape: BaseShape) -> str:
    """Serialise a pptx table shape to a GFM pipe table — same shape as
    the docx extractor produces."""
    if not shape.has_table:
        return ""
    table = shape.table  # type: ignore[attr-defined]
    rows = list(table.rows)
    if not rows:
        return ""

    def _cell_text(cell: Any) -> str:
        text = (cell.text or "").replace("\n", " ").replace("|", r"\|")
        return text.strip()

    header_cells = [_cell_text(c) for c in rows[0].cells]
    width = len(header_cells)
    if width == 0:
        return ""
    header_line = "| " + " | ".join(header_cells) + " |"
    sep_line = "|" + "|".join(["------"] * width) + "|"
    body_lines = []
    for row in rows[1:]:
        cells = [_cell_text(c) for c in row.cells]
        while len(cells) < width:
            cells.append("")
        body_lines.append("| " + " | ".join(cells[:width]) + " |")
    return "\n".join([header_line, sep_line, *body_lines])


def extract(path: Path) -> Iterator[Chunk]:
    reject_if_zip_bomb(path)
    try:
        yield from _extract_inner(path)
    except ExtractError:
        raise
    except Exception as e:
        # Parser libs raise everything from RuntimeError to opaque C
        # extension errors. Broad except is intentional — the goal is
        # "don't let one bad doc stop the index build."
        raise ExtractError(str(path), f"{type(e).__name__}: {e}") from e


def _extract_inner(path: Path) -> Iterator[Chunk]:
    parent_id = _parent_id(path)
    mtime = int(path.stat().st_mtime)
    try:
        prs = Presentation(str(path))
    except PackageNotFoundError as e:
        raise ExtractError(str(path), f"unreadable pptx: {e}") from e

    deck_title = ""
    for slide_index, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        if slide_index == 1 and title:
            deck_title = title

        body_parts: list[str] = []
        blocks: list[Block] = []
        md_lines: list[str] = []

        if title:
            blocks.append(Block(kind="h1", text=title))
            body_parts.append(title)
            md_lines.append(f"# {title}")

        title_shape = getattr(slide.shapes, "title", None)
        for shape in slide.shapes:
            if shape == title_shape:
                continue
            if shape.has_table:
                table_md = _table_md(shape)
                if table_md:
                    md_lines.append(table_md)
                    # Index every cell value so search lands on the
                    # right slide; one block per cell keeps body_struct
                    # plain-text for snippets.
                    for row in shape.table.rows:  # type: ignore[attr-defined]
                        for cell in row.cells:
                            cell_text = (cell.text or "").strip()
                            if cell_text:
                                blocks.append(Block(kind="p", text=cell_text))
                                body_parts.append(cell_text)
                continue
            shape_text = _shape_text(shape).strip()
            if not shape_text:
                continue
            md_lines.extend(_shape_md_lines(shape))
            blocks.append(Block(kind="p", text=shape_text))
            body_parts.append(shape_text)

        notes = _slide_notes(slide).strip()
        if notes:
            blocks.append(Block(kind="quote", text=notes))
            body_parts.append(notes)
            # Speaker notes render as a trailing blockquote; multi-line
            # notes need ``>`` on every line.
            note_lines = "\n".join(f"> {ln}" for ln in notes.splitlines() if ln.strip())
            if note_lines:
                md_lines.append(note_lines)

        body = "\n".join(body_parts).strip()
        if not body:
            continue

        body_md = "\n\n".join(line for line in md_lines if line).rstrip()

        yield Chunk(
            parent_id=parent_id,
            path=str(path),
            mtime=mtime,
            kind="pptx",
            body=body,
            body_struct=blocks,
            body_md=body_md,
            slide=slide_index,
            heading_path=title,
            title=deck_title,
            chunk_seq=slide_index - 1,
        )
