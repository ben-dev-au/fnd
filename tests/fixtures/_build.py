"""Build the test-fixture corpus.

Run with ``uv run python -m tests.fixtures._build``. Outputs are committed so the
contract is reproducible without running this script — but running it must produce
byte-identical content (no timestamps in PDF metadata, no random ordering).

Anchor phrases are deliberately weird so they cannot appear naturally elsewhere.
Each phrase must appear in exactly one (file, page/section) so retrieval tests can
assert exact attribution.
"""

from __future__ import annotations

import textwrap
from pathlib import Path

import pymupdf  # type: ignore[import-not-found]
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches, Pt

FIXTURES = Path(__file__).parent

# (path, page-or-section, content) — anchor phrases must be unique across the corpus.
# The generator inserts them so the test layer knows where to find them.
PDF_PAGES: dict[int, str] = {
    1: "Introduction. This is the first page of the test PDF document.",
    2: "Background. We discuss the relevant prior work in this section.",
    3: "Definitions. Key terms used throughout this document.",
    4: "Method. The general approach taken in our research.",
    5: "Implementation details that explain how things were built.",
    6: "Preliminary results from initial experiments.",
    7: "Detailed analysis. The blue penguin sandwich was observed only here.",
    8: "Discussion of limitations and future work.",
    9: "Related work in the broader academic literature.",
    10: "Acknowledgements to all collaborators.",
    11: "Bibliography and references.",
    12: "Appendix with supplementary material and proofs.",
}

MD_CONTENT = textwrap.dedent("""\
    # Test Notes

    This is the top-level body of the notes file.

    ## Methodology

    Notes about how the methodology section is structured.

    ### Sampling

    Details about sampling. The ostrich firewall was triggered during pilot runs.

    ### Analysis

    How analysis was performed.

    ## Results

    Summary of results.

    ## Conclusion

    Wrap-up paragraph.
    """)

TXT_CONTENT = textwrap.dedent("""\
    A short plain text file used for the TXT extractor tests.

    The marigold compiler is mentioned only here in the corpus; queries for it should
    surface this file at rank 1.

    A second paragraph keeps the file from being trivially small.
    """)


def _write_pdf(path: Path, pages: dict[int, str]) -> None:
    doc = pymupdf.open()
    try:
        for page_no in sorted(pages):
            page = doc.new_page(width=612, height=792)
            page.insert_text(
                (72, 100),
                f"Page {page_no}",
                fontsize=18,
                fontname="helv",
            )
            page.insert_text(
                (72, 150),
                pages[page_no],
                fontsize=12,
                fontname="helv",
            )
        path.parent.mkdir(parents=True, exist_ok=True)
        # garbage=4 + clean=True normalise output for reproducible bytes.
        doc.save(str(path), garbage=4, clean=True, deflate=True)
    finally:
        doc.close()


# ── Invisible-text fixture (scanned-OCR mimic) ───────────────────────────────
# Each page has a visible heading plus a body drawn in render-mode 3
# (invisible) — exactly how scanned books store their OCR layer behind a
# page image. pymupdf4llm's layout parser drops the invisible body
# (coverage collapses), so this fixture exercises the InvisibleTextTier
# recovery path deterministically without the real scanned corpus.
INVISIBLE_PROSE = (
    "The quicksort algorithm partitions recursively around a chosen pivot "
    "element producing sorted subsequences through comparison and exchange "
    "operations until the entire collection achieves total ordering."
)
# >20 distinct alpha tokens so the coverage gate's token floor is met.
INVISIBLE_CODE = [
    "class RingBuffer:",
    "    def __init__(self, capacity):",
    "        self.capacity = capacity",
    "        self.storage = list()",
    "    def enqueue(self, element):",
    "        self.storage.append(element)",
    "    def dequeue(self):",
    "        return self.storage.pop(first)",
    "    def is_empty(self):",
    "        return self.length() == zero",
    "    def length(self):",
    "        return len(self.storage)",
]


def _write_invisible_pdf(path: Path) -> None:
    doc = pymupdf.open()
    try:
        # Page 1: invisible prose behind a visible heading.
        page = doc.new_page(width=612, height=792)
        page.insert_text((72, 72), "Visible Heading Only", fontsize=14, fontname="helv")
        y = 120
        for line in (INVISIBLE_PROSE[i : i + 70] for i in range(0, len(INVISIBLE_PROSE), 70)):
            page.insert_text((72, y), line, fontsize=11, fontname="helv", render_mode=3)
            y += 16
        # Page 2: invisible monospace code (recovers as a ``` fence).
        page = doc.new_page(width=612, height=792)
        page.insert_text((72, 72), "Listing One", fontsize=13, fontname="helv")
        y = 110
        for line in INVISIBLE_CODE:
            page.insert_text((72, y), line, fontsize=10, fontname="cour", render_mode=3)
            y += 15
        path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(path), garbage=4, clean=True, deflate=True)
    finally:
        doc.close()


# ── Per-page heading fixture (doc-wide cutoff repro) ─────────────────────────
# pymupdf4llm builds one IdentifyHeaders over the whole document; with
# enough distinct large divider fonts the max_levels=6 cutoff lifts the
# body limit above genuine mid-size subheads, so they classify as body.
# Six invisible divider pages exhaust the cutoff; the target page's 16pt
# subhead is dropped doc-wide but recovered by per-page hdr_info.
HEADING_DIVIDER_SIZES = [30, 28, 26, 24, 22, 20]
HEADING_TARGET_BODY = (
    "Implementation notes describing the dispatcher routine which resolves "
    "handlers dynamically through registered factories binding concrete "
    "strategies lazily without explicit coupling between modules."
)


def _write_headings_pdf(path: Path) -> None:
    doc = pymupdf.open()
    try:
        for size in HEADING_DIVIDER_SIZES:
            page = doc.new_page(width=612, height=792)
            page.insert_text(
                (72, 72), "Section Divider", fontsize=size, fontname="helv", render_mode=3
            )
            page.insert_text(
                (72, 160),
                "filler body text paragraph alpha beta gamma delta",
                fontsize=11,
                fontname="helv",
                render_mode=3,
            )
        page = doc.new_page(width=612, height=792)
        page.insert_text(
            (72, 72), "Implementation Notes", fontsize=16, fontname="helv", render_mode=3
        )
        y = 160
        for line in (
            HEADING_TARGET_BODY[i : i + 70] for i in range(0, len(HEADING_TARGET_BODY), 70)
        ):
            page.insert_text((72, y), line, fontsize=11, fontname="helv", render_mode=3)
            y += 16
        # A visible running footer — real scanned pages carry a small
        # visible page label, so the layout path emits something (the gate
        # requires production Markdown to be present before it fires).
        page.insert_text((72, 760), "page seven", fontsize=9, fontname="helv")
        path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(path), garbage=4, clean=True, deflate=True)
    finally:
        doc.close()


def _write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


# ── PPTX fixture ─────────────────────────────────────────────────────────────
PPTX_SLIDES: list[tuple[str, str, str]] = [
    # (title, body_text, speaker_notes)
    ("Title Slide", "Test deck used by the PPTX extractor tests.", ""),
    ("Agenda", "Outline of the talk.", "Reminder: keep it short."),
    ("Background", "Some background information.", ""),
    (
        "Methods",
        "We discuss methods. The lavender stapler was the differentiating tool.",
        "This slide is the anchor; the unique phrase appears here only.",
    ),
    ("Results", "Outcomes of the pilot.", ""),
    ("Discussion", "What it means.", ""),
    ("Conclusion", "Wrap-up.", ""),
    ("Q&A", "Questions and answers.", ""),
]


def _write_pptx(path: Path, slides: list[tuple[str, str, str]]) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[5]  # title-only layout
    for title_text, body_text, notes_text in slides:
        slide = prs.slides.add_slide(blank)
        title_shape = slide.shapes.title
        if title_shape is not None:
            title_shape.text = title_text
        # Add a body text box.
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
        tf = tx_box.text_frame
        tf.text = body_text
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)
        if notes_text:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf is not None:
                notes_tf.text = notes_text
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


# ── DOCX fixture ─────────────────────────────────────────────────────────────
# A list of (style_name, text). Heading 1/2 establish heading_path; Body Text is body.
DOCX_PARAGRAPHS: list[tuple[str, str]] = [
    ("Heading 1", "Methods Document"),
    ("Body Text", "Top-level body of the document."),
    ("Heading 2", "Sampling"),
    ("Body Text", "Notes about sampling."),
    (
        "Body Text",
        "The narwhal compiler appears only in this section.",
    ),
    ("Heading 2", "Analysis"),
    ("Body Text", "Notes about analysis."),
    ("Heading 1", "Conclusion"),
    ("Body Text", "Wrap-up paragraph."),
]


def _write_docx(path: Path, paragraphs: list[tuple[str, str]]) -> None:
    doc = DocxDocument()
    for style, text in paragraphs:
        doc.add_paragraph(text, style=style)
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))


def build() -> None:
    _write_pdf(FIXTURES / "papers" / "test.pdf", PDF_PAGES)
    _write_invisible_pdf(FIXTURES / "scanned" / "invisible.pdf")
    _write_headings_pdf(FIXTURES / "scanned" / "headings.pdf")
    _write_text(FIXTURES / "notes" / "index.md", MD_CONTENT)
    _write_text(FIXTURES / "plain" / "short.txt", TXT_CONTENT)
    _write_pptx(FIXTURES / "slides" / "deck.pptx", PPTX_SLIDES)
    _write_docx(FIXTURES / "docs" / "methods.docx", DOCX_PARAGRAPHS)


# ── Anchor table — the test contract ─────────────────────────────────────────
# Each entry is (path-relative-to-FIXTURES, kind, locator, phrase).
# `locator` is a page number for PDF and a heading_path for MD.
ANCHORS: list[tuple[str, str, object, str]] = [
    ("papers/test.pdf", "pdf", 7, "blue penguin sandwich"),
    ("notes/index.md", "md", "Test Notes > Methodology > Sampling", "ostrich firewall"),
    ("plain/short.txt", "txt", None, "marigold compiler"),
    ("slides/deck.pptx", "pptx", 4, "lavender stapler"),
    ("docs/methods.docx", "docx", "Methods Document > Sampling", "narwhal compiler"),
]


if __name__ == "__main__":
    build()
    print("✓ fixtures built")
