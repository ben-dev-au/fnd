"""PPTX extractor emits markdown source on ``body_md`` for the new
structural preview renderer.

Slide title becomes ``# H1``; body shapes with non-zero-level
paragraphs render as bulleted lists with the level-derived indent;
flat-level shapes render as paragraphs; tables become GFM pipe
tables; speaker notes as a trailing blockquote.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from fnd.extract.pptx import extract


def _make_deck_with_bullets_and_table(out: Path) -> Path:
    prs = Presentation()
    blank = prs.slide_layouts[5]  # title-only

    # Slide 1: title + body shape with three bullets at increasing levels.
    s1 = prs.slides.add_slide(blank)
    s1.shapes.title.text = "Bullet Slide"  # type: ignore[union-attr]
    body = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(3))
    tf = body.text_frame
    tf.text = "outer alpha"
    tf.paragraphs[0].level = 0
    p2 = tf.add_paragraph()
    p2.text = "inner beta"
    p2.level = 1
    p3 = tf.add_paragraph()
    p3.text = "deeper gamma"
    p3.level = 2

    # Slide 2: title + a 2x2 table.
    s2 = prs.slides.add_slide(blank)
    s2.shapes.title.text = "Table Slide"  # type: ignore[union-attr]
    table = s2.shapes.add_table(
        rows=2, cols=2, left=Inches(1), top=Inches(2), width=Inches(6), height=Inches(2)
    ).table
    table.cell(0, 0).text = "Item"
    table.cell(0, 1).text = "Score"
    table.cell(1, 0).text = "alpha"
    table.cell(1, 1).text = "scaffolding"

    # Slide 3: speaker notes only — confirms blockquote path.
    s3 = prs.slides.add_slide(blank)
    s3.shapes.title.text = "Notes Slide"  # type: ignore[union-attr]
    s3.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2)).text_frame.text = "body text"
    s3.notes_slide.notes_text_frame.text = "remember to mention the dataset"  # type: ignore[union-attr]

    prs.save(str(out))
    return out


@pytest.fixture
def deck(tmp_path: Path) -> Path:
    return _make_deck_with_bullets_and_table(tmp_path / "fixture.pptx")


def test_slide_title_becomes_h1(deck: Path) -> None:
    chunks = list(extract(deck))
    assert chunks[0].body_md.lstrip().startswith("# Bullet Slide")
    assert chunks[1].body_md.lstrip().startswith("# Table Slide")


def test_bulleted_body_uses_list_indent_by_level(deck: Path) -> None:
    md = next(iter(extract(deck))).body_md
    assert "- outer alpha" in md
    assert "  - inner beta" in md
    assert "    - deeper gamma" in md


def test_table_renders_as_gfm_pipe_table(deck: Path) -> None:
    md = list(extract(deck))[1].body_md
    assert "| Item | Score |" in md
    assert "| alpha | scaffolding |" in md
    assert "|------|" in md.replace(" ", "")


def test_table_cell_text_is_indexed(deck: Path) -> None:
    """Cell values must land in chunk.body so F_BODY can search them."""
    chunks = list(extract(deck))
    assert "scaffolding" in chunks[1].body, chunks[1].body


def test_speaker_notes_become_trailing_blockquote(deck: Path) -> None:
    md = list(extract(deck))[2].body_md
    assert "> remember to mention the dataset" in md, md
    # And it lands AFTER the slide body, not before.
    assert md.index("body text") < md.index("> remember")


def test_per_slide_chunk_seq_unchanged(deck: Path) -> None:
    """Slide indexing must keep sequencing from 0 in document order."""
    chunks = list(extract(deck))
    assert [c.chunk_seq for c in chunks] == [0, 1, 2]
    assert [c.slide for c in chunks] == [1, 2, 3]


def test_body_struct_remains_plain_text(deck: Path) -> None:
    """Snippet pipeline reads body_struct — markdown markers in body_md
    must not leak into block.text."""
    for chunk in extract(deck):
        for block in chunk.body_struct:
            assert "**" not in block.text, block.text


def test_no_change_to_existing_simple_deck_extraction(tmp_path: Path) -> None:
    """A deck that's just title + flat-level body (the legacy fixture
    shape) keeps producing the same searchable body — no regression on
    existing test_pptx_anchor_attributes_to_correct_slide."""
    prs = Presentation()
    blank = prs.slide_layouts[5]
    s = prs.slides.add_slide(blank)
    s.shapes.title.text = "Methods"  # type: ignore[union-attr]
    body = s.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(3))
    body.text_frame.text = "We discuss methods. The lavender stapler was the differentiating tool."
    out = tmp_path / "simple.pptx"
    prs.save(str(out))
    chunks = list(extract(out))
    assert len(chunks) == 1
    assert "lavender stapler" in chunks[0].body
    # Flat (level=0) body is rendered as a paragraph, not a bullet item.
    assert "- We discuss" not in chunks[0].body_md
    assert "We discuss methods." in chunks[0].body_md
