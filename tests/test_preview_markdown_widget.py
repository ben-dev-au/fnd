"""Structural markdown rendering in the preview pane.

Replaces the legacy ``test_md_match_chunk_uses_per_line_layout`` —
matched markdown chunks now mount through ``FNDMarkdown`` (Textual's
Markdown widget with highlight-aware block subclasses) so tables, code
fences, lists, and blockquotes render structurally even when they
contain query matches. The highlight overlay is a per-span mark inside
the rendered Content, NOT a per-line widget swap.
"""

from __future__ import annotations

import textwrap
from pathlib import Path

import pytest
from textual.containers import VerticalScroll
from textual.widgets import Tree
from textual.widgets._markdown import (
    MarkdownFence,
    MarkdownTable,
)

from fnd.config import Config, load
from fnd.index import build_index
from fnd.render import HIGHLIGHT_STYLE
from fnd.tui import FNDApp
from fnd.tui.app import (
    FNDMarkdown,
    FNDMarkdownParagraph,
)


def _is_highlight_span(span: object) -> bool:
    """The renderer sets the search-highlight Span style to the same
    Rich style string the per-line plain renderer uses
    (``fnd.render.HIGHLIGHT_STYLE``). The Span style attribute is
    parsed into a Style instance, so compare via the string form."""
    style = getattr(span, "style", None)
    return str(style) == HIGHLIGHT_STYLE or HIGHLIGHT_STYLE in str(style)


def _write(p: Path, body: str) -> None:
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(body, encoding="utf-8")


@pytest.fixture
def cfg(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Config:
    cfg_path = tmp_path / "config.toml"
    cfg_path.write_text(
        textwrap.dedent("""
            [[collections.notes.sources]]
            path = "/tmp/notes"
        """),
        encoding="utf-8",
    )
    monkeypatch.setattr("fnd.config.default_config_path", lambda: cfg_path)
    return load(cfg_path)


# ── Highlights are match-only, not paragraph-wide ────────────────────


@pytest.fixture
def paragraph_index(tmp_path: Path, tmp_index_dir: Path) -> Path:
    a = tmp_path / "notes"
    _write(
        a / "Notes.md",
        textwrap.dedent(
            """\
            # Patterns

            The templates pattern is described on page 47.
            """
        ),
    )
    build_index(roots=[a], index_dir=tmp_index_dir, collection="notes")
    return tmp_index_dir


@pytest.mark.asyncio
async def test_preview_md_match_only_highlights_term(cfg: Config, paragraph_index: Path) -> None:
    """A query that matches one word in a paragraph applies a span over
    exactly that word's character range — not the whole paragraph."""
    app = FNDApp(
        index_dir=paragraph_index,
        config=cfg,
        collection="notes",
        initial_query="templates",
    )
    async with app.run_test() as pilot:
        await pilot.pause()
        tree = app.query_one("#results_pane", Tree)
        tree.focus()
        await pilot.pause()
        pane = app.query_one("#preview_pane", VerticalScroll)
        # Exactly one FNDMarkdown widget mounted (one chunk, one paragraph).
        md_widgets = list(pane.query(FNDMarkdown))
        assert len(md_widgets) == 1
        # The paragraph block carries highlight spans.
        paras = list(pane.query(FNDMarkdownParagraph))
        assert paras
        spans = list(paras[0]._content.spans)
        # A span exists, has the search-highlight style, and bounds the
        # term "templates" exactly.
        highlight_spans = [s for s in spans if _is_highlight_span(s)]
        assert len(highlight_spans) == 1, spans
        plain = paras[0]._content.plain
        s = highlight_spans[0]
        assert plain[s.start : s.end].lower() == "templates"
        # Stem-aware: a slightly different inflection would still match,
        # but here we asserted exact word so the bounds must be exact.


# ── Tables render with per-cell highlights, not collapsed ─────────────


@pytest.fixture
def table_index(tmp_path: Path, tmp_index_dir: Path) -> Path:
    a = tmp_path / "notes"
    _write(
        a / "Tables.md",
        textwrap.dedent(
            """\
            # Costs

            | Item | Price |
            |------|-------|
            | apple | 1.20 |
            | beanstalk | 99.00 |
            """
        ),
    )
    build_index(roots=[a], index_dir=tmp_index_dir, collection="notes")
    return tmp_index_dir


@pytest.mark.asyncio
async def test_preview_md_renders_table_with_cell_highlight(cfg: Config, table_index: Path) -> None:
    """Tables render as actual tables; a query matching one cell value
    highlights only that cell. The table compose flow folds the
    MarkdownTH/MarkdownTD widgets' _content (with our spans already
    applied) into the MarkdownTableContent grid, so we read the spans
    off ``_headers`` / ``_rows`` directly rather than querying TH/TD
    widgets — the latter aren't mounted as visible widgets, only
    referenced as Content sources."""
    app = FNDApp(
        index_dir=table_index,
        config=cfg,
        collection="notes",
        initial_query="beanstalk",
    )
    async with app.run_test() as pilot:
        await pilot.pause()
        tree = app.query_one("#results_pane", Tree)
        tree.focus()
        await pilot.pause()
        pane = app.query_one("#preview_pane", VerticalScroll)
        tables = list(pane.query(MarkdownTable))
        assert tables, "table chunk should render via MarkdownTable widget"
        table = tables[0]
        # Header content carries no highlights (query doesn't match
        # "Item"/"Price").
        for header_content in table._headers:
            assert all(
                s.style != "search-highlight" for s in header_content.spans
            ), header_content.spans
        # Body rows: exactly one cell across all rows has the highlight,
        # and it bounds the matched substring.
        all_cell_contents = [cell for row in table._rows for cell in row]
        highlighted = [
            cell for cell in all_cell_contents if any(_is_highlight_span(s) for s in cell.spans)
        ]
        assert (
            len(highlighted) == 1
        ), f"expected exactly one highlighted cell, got {len(highlighted)}"
        cell = highlighted[0]
        spans = [s for s in cell.spans if _is_highlight_span(s)]
        assert len(spans) == 1
        assert cell.plain[spans[0].start : spans[0].end].lower() == "beanstalk"


# ── Code fences are not highlighted ──────────────────────────────────


@pytest.fixture
def fence_index(tmp_path: Path, tmp_index_dir: Path) -> Path:
    """Heading carries the searchable term so the fusion path returns
    this chunk; the fenced code block separately contains the term so
    the highlight subclass would (incorrectly) wrap it if it were
    registered for fences. The test asserts that doesn't happen."""
    a = tmp_path / "notes"
    _write(
        a / "Code.md",
        textwrap.dedent(
            """\
            # Templates Snippet

            A templates example follows below.

            ```python
            templates = ["a", "b"]
            ```
            """
        ),
    )
    build_index(roots=[a], index_dir=tmp_index_dir, collection="notes")
    return tmp_index_dir


@pytest.mark.asyncio
async def test_preview_md_fence_no_highlight_inside_code(cfg: Config, fence_index: Path) -> None:
    """Query terms inside a fenced code block must NOT receive the
    highlight overlay — code stays on Textual's stock MarkdownFence
    rendering (rich.syntax.Syntax) so syntax colours remain readable.
    """
    app = FNDApp(
        index_dir=fence_index,
        config=cfg,
        collection="notes",
        initial_query="templates",
    )
    async with app.run_test() as pilot:
        await pilot.pause()
        tree = app.query_one("#results_pane", Tree)
        tree.focus()
        await pilot.pause()
        pane = app.query_one("#preview_pane", VerticalScroll)
        fences = list(pane.query(MarkdownFence))
        assert fences, "code-only chunk should render via MarkdownFence"
        # MarkdownFence isn't subclassed by us → its content carries no
        # search-highlight spans regardless of what the query is.
        for fence in fences:
            spans = list(fence._content.spans)
            assert all(s.style != "search-highlight" for s in spans), spans


# ── Nested lists render structurally ─────────────────────────────────


@pytest.fixture
def nested_list_index(tmp_path: Path, tmp_index_dir: Path) -> Path:
    a = tmp_path / "notes"
    _write(
        a / "Lists.md",
        textwrap.dedent(
            """\
            # Outline

            - outer one
              - inner alpha
              - inner beta
            - outer two
            """
        ),
    )
    build_index(roots=[a], index_dir=tmp_index_dir, collection="notes")
    return tmp_index_dir


@pytest.mark.asyncio
async def test_preview_md_nested_lists_render(cfg: Config, nested_list_index: Path) -> None:
    """Nested bullets render as nested list widgets (multiple
    MarkdownBulletList ancestors deep), not as flat paragraphs."""
    from textual.widgets._markdown import MarkdownBulletList

    app = FNDApp(
        index_dir=nested_list_index,
        config=cfg,
        collection="notes",
        initial_query="alpha",
    )
    async with app.run_test() as pilot:
        await pilot.pause()
        tree = app.query_one("#results_pane", Tree)
        tree.focus()
        await pilot.pause()
        pane = app.query_one("#preview_pane", VerticalScroll)
        bullet_lists = list(pane.query(MarkdownBulletList))
        # Outer + inner list = at least two MarkdownBulletList widgets.
        assert len(bullet_lists) >= 2, f"expected nested bullet lists; got {len(bullet_lists)}"


# ── first_match_block ─────────────────────────────────────────────────


@pytest.fixture
def multi_para_index(tmp_path: Path, tmp_index_dir: Path) -> Path:
    a = tmp_path / "notes"
    _write(
        a / "Paragraphs.md",
        textwrap.dedent(
            """\
            # Many

            paragraph one is here.

            paragraph two has nothing.

            paragraph three with templates landing here.

            paragraph four trailing.
            """
        ),
    )
    build_index(roots=[a], index_dir=tmp_index_dir, collection="notes")
    return tmp_index_dir


@pytest.fixture
def docx_corpus(tmp_path: Path, tmp_index_dir: Path) -> Path:
    """A docx with a heading + bold + a table — exercises the docx
    extractor's body_md path end-to-end."""
    from docx import Document

    a = tmp_path / "docs"
    a.mkdir(parents=True)
    doc = Document()
    doc.add_heading("Paperwork", level=1)
    p = doc.add_paragraph()
    p.add_run("plain ")
    p.add_run("templates").bold = True
    p.add_run(" appear here.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Category"
    table.cell(0, 1).text = "Notes"
    table.cell(1, 0).text = "Type"
    table.cell(1, 1).text = "scaffolding"
    doc.save(str(a / "Patterns.docx"))
    build_index(roots=[a], index_dir=tmp_index_dir, collection="notes")
    return tmp_index_dir


@pytest.fixture
def pptx_corpus(tmp_path: Path, tmp_index_dir: Path) -> Path:
    """Single-slide deck with a title, bulleted body, and a table —
    exercises the pptx extractor's body_md path end-to-end."""
    from pptx import Presentation
    from pptx.util import Inches

    a = tmp_path / "decks"
    a.mkdir(parents=True)
    prs = Presentation()
    blank = prs.slide_layouts[5]
    s1 = prs.slides.add_slide(blank)
    s1.shapes.title.text = "Strategy"  # type: ignore[union-attr]
    body = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(3))
    tf = body.text_frame
    tf.text = "templates pattern overview"
    tf.paragraphs[0].level = 0
    p = tf.add_paragraph()
    p.text = "iterators"
    p.level = 1
    table = s1.shapes.add_table(
        rows=2, cols=2, left=Inches(1), top=Inches(5), width=Inches(6), height=Inches(2)
    ).table
    table.cell(0, 0).text = "Pattern"
    table.cell(0, 1).text = "Use"
    table.cell(1, 0).text = "Visitor"
    table.cell(1, 1).text = "templates"
    prs.save(str(a / "Patterns.pptx"))
    build_index(roots=[a], index_dir=tmp_index_dir, collection="notes")
    return tmp_index_dir


@pytest.mark.asyncio
async def test_pptx_preview_routes_through_fnd_markdown(cfg: Config, pptx_corpus: Path) -> None:
    from textual.widgets._markdown import MarkdownTable

    app = FNDApp(
        index_dir=pptx_corpus,
        config=cfg,
        collection="notes",
        initial_query="templates",
    )
    async with app.run_test() as pilot:
        await pilot.pause()
        tree = app.query_one("#results_pane", Tree)
        tree.focus()
        await pilot.pause()
        pane = app.query_one("#preview_pane", VerticalScroll)
        assert list(pane.query(FNDMarkdown)), "pptx chunk should mount FNDMarkdown"
        assert list(pane.query(MarkdownTable)), "pptx table should render via MarkdownTable"


@pytest.mark.asyncio
async def test_docx_preview_routes_through_fnd_markdown(cfg: Config, docx_corpus: Path) -> None:
    """A docx chunk's preview mounts an FNDMarkdown widget and
    renders the embedded table — proves the renderer dispatch wires
    the docx kind through the new path."""
    from textual.widgets._markdown import MarkdownTable

    app = FNDApp(
        index_dir=docx_corpus,
        config=cfg,
        collection="notes",
        initial_query="templates",
    )
    async with app.run_test() as pilot:
        await pilot.pause()
        tree = app.query_one("#results_pane", Tree)
        tree.focus()
        await pilot.pause()
        pane = app.query_one("#preview_pane", VerticalScroll)
        assert list(pane.query(FNDMarkdown)), "docx chunk should mount FNDMarkdown"
        assert list(pane.query(MarkdownTable)), "docx table should render via MarkdownTable"


@pytest.mark.asyncio
async def test_preview_first_match_block_resolves_to_matched_paragraph(
    cfg: Config, multi_para_index: Path
) -> None:
    """``FNDMarkdown.first_match_block`` should resolve to the third
    paragraph (the one containing the query term), not paragraph one."""
    app = FNDApp(
        index_dir=multi_para_index,
        config=cfg,
        collection="notes",
        initial_query="templates",
    )
    async with app.run_test() as pilot:
        await pilot.pause()
        tree = app.query_one("#results_pane", Tree)
        tree.focus()
        await pilot.pause()
        pane = app.query_one("#preview_pane", VerticalScroll)
        md = pane.query_one(FNDMarkdown)
        first_match = md.first_match_block  # type: ignore[attr-defined]
        assert first_match is not None
        assert "templates" in first_match._content.plain
